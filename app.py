import gradio as gr
import os
import tempfile
import zipfile
import shutil
import base64
import json
import re
import concurrent.futures
import time
from pdf2image import convert_from_path
from PIL import Image
from dotenv import load_dotenv

# PPTX 處理套件
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 使用 Google 新版 SDK
from google import genai
from google.genai import types

load_dotenv()

class NotebookLMTool:
    def __init__(self):
        self.api_key = os.getenv("GEMINI_API_KEY")
        self.client = None
        if self.api_key:
            self.client = genai.Client(api_key=self.api_key)

    def set_key(self, user_key):
        if user_key and user_key.strip():
            self.api_key = user_key.strip()
            self.client = genai.Client(api_key=self.api_key)
            return "✅ API Key 已更新！"
        return "⚠️ Key 無效"

    def _extract_json(self, text):
        """強化版 JSON 提取"""
        try:
            match = re.search(r"```json\s*(.*)\s*```", text, re.DOTALL)
            if match: return json.loads(match.group(1))
            match = re.search(r"\[\s*\{.*\}\s*\]", text, re.DOTALL)
            if match: return json.loads(match.group(0))
            return json.loads(text)
        except:
            return []

    # --- 單頁處理邏輯 (獨立出來以便平行運算) ---
    def process_single_page(self, page_index, img, img_output_dir):
        """處理單一頁面的：去字(背景) + 文字分析(Layout)"""
        print(f"🚀 [Page {page_index+1}] 開始處理...", flush=True)
        
        # 結果容器
        result = {
            "index": page_index,
            "bg_path": None,
            "blocks": [],
            "log": "",
            "preview": None,
            "tokens_in": 0,
            "tokens_out": 0
        }

        save_name = f"slide_{page_index+1:02d}.png"
        final_bg_path = os.path.join(img_output_dir, save_name)
        bg_success = False

        # 1. 背景去字 (Image Cleaning)
        try:
            clean_prompt = """
            Strictly remove all text, titles, text-boxes, and bullet points from this slide image.
            CRITICAL INSTRUCTION:
            1. Preserve the original background pattern, colors, logos, and non-text graphics EXACTLY as they are.
            2. Do NOT add any new objects, decorations, or hallucinations.
            3. Output ONLY the image.
            """
            
            # 使用 2.0-flash-exp 進行繪圖
            resp_img = self.client.models.generate_content(
                model="gemini-2.5-flash-image", 
                contents=[clean_prompt, img],
                config=types.GenerateContentConfig(response_modalities=["IMAGE"])
            )

            # Token 統計
            if resp_img.usage_metadata:
                result["tokens_in"] += resp_img.usage_metadata.prompt_token_count
                result["tokens_out"] += resp_img.usage_metadata.candidates_token_count

            # 存圖邏輯
            image_data = None
            if hasattr(resp_img, 'parts') and resp_img.parts:
                for part in resp_img.parts:
                    if part.inline_data: image_data = part.inline_data.data; break
            if image_data is None and hasattr(resp_img, 'bytes') and resp_img.bytes:
                image_data = resp_img.bytes

            if image_data:
                if isinstance(image_data, str): image_data = base64.b64decode(image_data)
                with open(final_bg_path, "wb") as f: f.write(image_data)
                bg_success = True
                result["bg_path"] = final_bg_path
                result["preview"] = (final_bg_path, f"Page {page_index+1} Cleaned")
            else:
                print(f"⚠️ [Page {page_index+1}] 去字失敗: 未回傳圖片", flush=True)

        except Exception as e:
            print(f"❌ [Page {page_index+1}] Clean Error: {e}", flush=True)

        # 失敗回退原圖
        if not bg_success:
            img.save(final_bg_path)
            result["bg_path"] = final_bg_path # 仍需路徑給 PPT 使用
            result["preview"] = (final_bg_path, f"Page {page_index+1} (Original)")
            result["log"] += f"[P{page_index+1}] Warning: Background cleaning failed.\n"

        # 2. 文字與佈局分析 (Layout Analysis)
        try:
            layout_prompt = """
            Analyze this slide. Return a JSON list of all text blocks.
            Each item: {"text": string, "box_2d": [ymin, xmin, ymax, xmax] (0-1000), "font_size": int, "color": hex, "is_bold": bool}
            """
            
            resp_layout = self.client.models.generate_content(
                model="gemini-2.5-flash", 
                contents=[layout_prompt, img],
                config=types.GenerateContentConfig(response_mime_type="application/json")
            )

            if resp_layout.usage_metadata:
                result["tokens_in"] += resp_layout.usage_metadata.prompt_token_count
                result["tokens_out"] += resp_layout.usage_metadata.candidates_token_count

            blocks = self._extract_json(resp_layout.text)
            result["blocks"] = blocks
            
            # 紀錄 Log
            for b in blocks:
                if b.get("text"): result["log"] += f"[P{page_index+1}] {b['text'][:20]}...\n"

        except Exception as e:
            print(f"❌ [Page {page_index+1}] Layout Error: {e}", flush=True)
            result["log"] += f"[P{page_index+1}] Layout Analysis Failed.\n"

        print(f"✅ [Page {page_index+1}] 完成！", flush=True)
        return result

    def process_pdf(self, pdf_file, progress=gr.Progress()):
        if not self.client:
            raise ValueError("請先輸入 Google API Key！")
        
        if pdf_file is None:
            return None, None, None, ""

        # 統計數據
        total_input_tokens = 0
        total_output_tokens = 0
        full_text_log = ""
        gallery_preview = []
        
        # 1. 準備環境
        temp_dir = tempfile.mkdtemp()
        img_output_dir = os.path.join(temp_dir, "cleaned_images")
        os.makedirs(img_output_dir, exist_ok=True)
        
        # 初始化 PPTX
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        # 2. PDF 轉圖片 (降低 DPI 加速)
        progress(0.1, desc="正在將 PDF 轉為圖片 (DPI=150)...")
        try:
            # dpi=150 足夠螢幕檢視，且大幅減少上傳時間
            images = convert_from_path(pdf_file, dpi=150) 
        except Exception as e:
            raise ValueError(f"PDF 轉換失敗: {str(e)}")

        # 3. 平行處理 (Parallel Execution)
        # 根據 CPU 核心數或 API 限制設定 workers，建議 3-5 避免 Rate Limit
        max_workers = 4 
        results_map = {} # 用來存結果，確保順序正確

        progress(0.2, desc="🚀 AI 多工處理中 (可能需要稍等)...")
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            # 提交所有任務
            future_to_page = {
                executor.submit(self.process_single_page, i, img, img_output_dir): i 
                for i, img in enumerate(images)
            }
            
            # 等待完成
            for future in concurrent.futures.as_completed(future_to_page):
                try:
                    res = future.result()
                    results_map[res["index"]] = res
                    # 更新 Token
                    total_input_tokens += res["tokens_in"]
                    total_output_tokens += res["tokens_out"]
                except Exception as exc:
                    print(f"Page processing generated an exception: {exc}")

        # 4. 依序組裝 PPTX (確保順序正確)
        progress(0.8, desc="正在組裝 PPTX...")
        
        cleaned_images_paths = [] # 用於 ZIP
        
        for i in range(len(images)):
            if i not in results_map:
                print(f"Missing result for page {i}")
                continue
                
            res = results_map[i]
            
            # 更新 Log 與 Preview
            full_text_log += res["log"]
            if res["preview"]: gallery_preview.append(res["preview"])
            if res["bg_path"]: cleaned_images_paths.append(res["bg_path"])

            # 建立 Slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # A. 貼背景
            if res["bg_path"] and os.path.exists(res["bg_path"]):
                try:
                    slide.shapes.add_picture(res["bg_path"], 0, 0, width=prs.slide_width, height=prs.slide_height)
                except: pass
            
            # B. 貼文字
            for block in res["blocks"]:
                text_content = block.get("text", "")
                if not text_content: continue
                
                # 座標轉換
                box = block.get("box_2d", [0, 0, 100, 100])
                ymin, xmin, ymax, xmax = box
                left = Inches((xmin / 1000) * 16)
                top = Inches((ymin / 1000) * 9)
                width = Inches(((xmax - xmin) / 1000) * 16)
                height = Inches(((ymax - ymin) / 1000) * 9)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text_content
                
                try: p.font.size = Pt(int(block.get("font_size", 18)))
                except: p.font.size = Pt(18)
                p.font.bold = block.get("is_bold", False)
                try:
                    hex_c = block.get("color", "#000000").replace("#", "")
                    # 如果背景去字失敗，原圖背景可能很複雜，文字顏色可能需要調整 (這裡暫不處理，保持原色)
                    p.font.color.rgb = RGBColor.from_string(hex_c)
                except: pass

        # 5. 打包
        progress(0.9, desc="正在打包檔案...")
        pptx_path = os.path.join(temp_dir, "restored_presentation.pptx")
        prs.save(pptx_path)

        txt_path = os.path.join(temp_dir, "content_log.txt")
        with open(txt_path, "w", encoding="utf-8") as f: f.write(full_text_log)

        zip_path = os.path.join(temp_dir, "notebooklm_restore_pack.zip")
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            zf.write(pptx_path, "restored_slides.pptx")
            zf.write(txt_path, "content_log.txt")
            for img_path in cleaned_images_paths:
                zf.write(img_path, os.path.join("cleaned_backgrounds", os.path.basename(img_path)))

        token_stats = f"""
        ### 📊 Token 用量統計
        - **總輸入:** {total_input_tokens:,}
        - **總輸出:** {total_output_tokens:,}
        - **總計消耗:** {total_input_tokens + total_output_tokens:,}
        """

        return zip_path, pptx_path, gallery_preview, token_stats

# Init
tool = NotebookLMTool()

# --- Gradio UI ---
with gr.Blocks(title="NotebookLM Slide Restorer，PPT.404", theme=gr.themes.Soft()) as demo:
    gr.Markdown("# 🛠️ NotebookLM 投影片 PDF 還原神器 (PPT.404)")
    gr.Markdown("""
    <div align="center">
    
    # 🪄 上傳 PDF，AI 自動：**去字背景** + **版面分析** + **合成可編輯 PPTX**
    👉 歡迎 Star [GitHub](https://github.com/Deep-Learning-101/) ⭐ 覺得不錯 👈  
    
    <h3>🧠 補腦專區：<a href="https://deep-learning-101.github.io/" target="_blank">Deep Learning 101</a></h3>  
    
    | 🔥 技術傳送門 (Tech Stack) | 📚 必讀心法 (Must Read) |
    | :--- | :--- |
    | 🤖 [**大語言模型 (LLM)**](https://deep-learning-101.github.io/Large-Language-Model) | 🏹 [**策略篇：企業入門策略**](https://deep-learning-101.github.io/Blog/AIBeginner) |
    | 📝 [**自然語言處理 (NLP)**](https://deep-learning-101.github.io/Natural-Language-Processing) | 📊 [**評測篇：臺灣 LLM 分析**](https://deep-learning-101.github.io/Blog/TW-LLM-Benchmark) |
    | 👁️ [**電腦視覺 (CV)**](https://deep-learning-101.github.io//Computer-Vision) | 🛠️ [**實戰篇：打造高精準 RAG**](https://deep-learning-101.github.io/RAG) |
    | 🎤 [**語音處理 (Speech)**](https://deep-learning-101.github.io/Speech-Processing) | 🕳️ [**避坑篇：AI Agent 開發陷阱**](https://deep-learning-101.github.io/agent) |
    </div>
    """)
    
    with gr.Row():
        with gr.Column():
            api_input = gr.Textbox(label="Google API Key", type="password", placeholder="貼上你的 Gemini API Key")
            btn_set_key = gr.Button("設定 Key")
            status_msg = gr.Markdown("")
            
            gr.Markdown("---")
            pdf_input = gr.File(label="上傳 PDF")
            btn_process = gr.Button("🚀 開始還原 PPTX (平行加速版)", variant="primary")
        
        with gr.Column():
            out_zip = gr.File(label="📦 下載完整包")
            out_pptx = gr.File(label="📊 直接下載 PPTX")
            out_tokens = gr.Markdown("### 📊 等待處理...")
    
    gr.Markdown("### 🖼️ 背景去字效果預覽")
    out_gallery = gr.Gallery(columns=4)

    btn_set_key.click(tool.set_key, inputs=api_input, outputs=status_msg)
    
    btn_process.click(
        tool.process_pdf,
        inputs=[pdf_input],
        outputs=[out_zip, out_pptx, out_gallery, out_tokens]
    )

if __name__ == "__main__":
    demo.launch()
